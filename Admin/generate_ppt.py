"""
Generate PowerPoint presentation from markdown slides
CRP Second Debate Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "From Static to Dynamic: Parameter Manipulation as Adaptive DDL Scaffolding in Creative Writing with Generative AI"
    subtitle.text = "Ruobin Yu\nCRP Second Debate\nHong Kong Baptist University\n\nRevised Draft: Addressing Panel Feedback"
    
    # Slide 2: Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Agenda - Key Revisions Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Six Major Changes in Response to Panel Feedback:"
    
    points = [
        "1. Narrowed Research Focus → Single unified RQ on parameter effects",
        "2. L2 Framing Reconsidered → 'Novice writers' emphasis + limitations acknowledged",
        "3. Poetry vs. Machine Logic → Heidegger challenged, writing-as-problem-solving defended",
        "4. DDL Justification → AI as 'new DDL interface' argument developed",
        "5. Creativity Framework → Type C and divergent creativity mechanism explained",
        "6. Limitations & Future Studies → Comparison groups, longitudinal design proposed"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
    
    # Slide 3: Change #1 - Narrowed Research Focus
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #1 - Narrowed Research Focus"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "From Broad to Focused"
    
    p = tf.add_paragraph()
    p.text = "Original Issue (Dr. Harrington):"
    p.font.bold = True
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = '"Research questions are too broad"'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Revision:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "RQ1: How do different parameter configurations influence authorship/ownership, creative agency, and satisfaction for novice poetry writers in university classroom?"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "RQ2: What pedagogical implications emerge for configuring AI as adaptive DDL scaffolding?"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Key Refinement:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    refinements = [
        "Unified focus on parameter → interaction → perception mechanism",
        "Explicit outcome measures: authorship, agency, satisfaction",
        "Context: university creative writing classrooms, novice writers"
    ]
    
    for ref in refinements:
        p = tf.add_paragraph()
        p.text = ref
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 4: Change #2 - L2 Framing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #2 - L2 Framing Reconsidered"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Avoiding Deficit Assumptions"
    
    p = tf.add_paragraph()
    p.text = "Original Issue (Dr. Checketts, Dr. Hung):"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Avoid assuming L2 learners are deficient compared to L1 writers"'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Revision Strategy:"
    p.font.bold = True
    
    strategies = [
        "✓ Emphasis on 'novice writers' rather than L2/L1 distinction",
        "✓ Limitation section acknowledges lack of L1 comparison group",
        "✓ Future studies propose L2 and L1 parallel groups"
    ]
    
    for strat in strategies:
        p = tf.add_paragraph()
        p.text = strat
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Rationale: Focus on pedagogical scaffolding needs rather than language proficiency deficits"
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Slide 5: Change #3 - Poetry vs. Machine Logic
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #3 - Poetry vs. Machine Logic"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Challenging Heidegger's Binary"
    
    p = tf.add_paragraph()
    p.text = "Heidegger (1984) vs. Strehovec (2023)"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Heidegger: Poetry = purest human expression, incompatible with machine logic"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Strehovec: Artistic creation = cognitive activity, poetry = problem-solving"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Our Position:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Machine influence is pedagogically valid when bounded and transparent: AI provides pattern data and options, while humans retain selection, rejection, and revision authority."
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "For beginners in writing who have ideas but lack linguistic resources, AI-assisted poetry can be DDL when the process foregrounds noticing, testing, and reflective decision making."
    p.level = 1
    p.font.size = Pt(13)
    
    # Slide 6: Change #4 - DDL Justification
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #4 - DDL Justification"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "AI as a New DDL Interface"
    
    p = tf.add_paragraph()
    p.text = "Traditional DDL:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Learner → Query → Corpus → Concordance → Discovery"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "AI-Assisted DDL:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Learner + AI → Pattern Generation → Noticing/Analysis → Selection/Revision"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Key Statement:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"AI can serve as a new DDL interface if learners treat its outputs as linguistic data to notice, analyze, and revise."'
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Why Poetry? Traditional DDL shows 'what has been written'; AI-assisted DDL enables 'what could be written' → critical for creative writing where novelty and surprise are pedagogical goals"
    p.level = 1
    p.font.size = Pt(13)
    
    # Slide 7: Parameter → Interaction → Perception Mechanism
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Parameter → Interaction → Perception Mechanism"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Technical-Pedagogical Bridge"
    
    p = tf.add_paragraph()
    p.text = "The Mechanism We Discovered:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Parameter Configuration ↓"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Interaction Type Distribution ↓"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Authorship & Agency Perception"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Empirical Evidence:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Structured (T=0.3, top-p=0.4): Type C 5%, Lower authorship (40-60%)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Exploratory (T=0.8, top-p=0.9): Type C 35%, Higher authorship (60-96%)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "7-fold increase in Type C (χ²=24.3, p<.001, Cramer's V=.38)"
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    
    # Slide 8: Three Interaction Types Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Three Interaction Types Framework"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "DDL-Scaffolding Modes"
    
    p = tf.add_paragraph()
    p.text = "Type A: Constraint Repair"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Pattern reinforcement | Low parameter | Heavy scaffolding | Convergent creativity"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Type B: Exemplar Giving"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Model selection | Medium parameter | Moderate scaffolding | Hybrid creativity"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Type C: Surprise Harvest"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Serendipitous discovery | High parameter | Light scaffolding | Divergent creativity"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = 'Type A: "AI as tutor" | Type B: "AI as resource bank" | Type C: "AI as creative partner"'
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Slide 9: The Type B Paradox
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Type B Paradox"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Helpful but Alienating"
    
    p = tf.add_paragraph()
    p.text = "Empirical Finding:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✓ Most common (61.54% voted)"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "✓ Most helpful (55.56% voted)"
    p.level = 1
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "✗ Reduced authorship perception"
    p.level = 1
    p.font.color.rgb = RGBColor(192, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "✗ Lower creative satisfaction"
    p.level = 1
    p.font.color.rgb = RGBColor(192, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "Participant Quotes:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = 'P02: "I do not feel ownership or pride over much of it"'
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = 'P05: "Feels like just AI... I couldn\'t feel any sort of creativity"'
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Cognitive Shift: From 'How do I create?' to 'Which option do I choose?'"
    p.font.italic = True
    p.font.color.rgb = RGBColor(0, 0, 128)
    
    # Slide 10: Change #5 - Creativity Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #5 - Creativity Framework in Discussion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Type C: Convergent → Divergent Transformation"
    
    p = tf.add_paragraph()
    p.text = "Type C Mechanism:"
    p.font.bold = True
    
    mechanisms = [
        "1. AI generates unexpected suggestions (divergent input)",
        "2. Human exercises critical judgment (evaluation)",
        "3. Iterative reciprocity emerges (negotiation cycles)",
        "4. Authorship activated through rejection/refinement (agency)"
    ]
    
    for mech in mechanisms:
        p = tf.add_paragraph()
        p.text = mech
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Evidence:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• P09: 80% authorship through 5 iterative cycles"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• P12: 90% authorship (highest across 30 participants)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• P16: 50-50 co-creation with maximum satisfaction (5/5)"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Core Insight: Type C provokes divergent thinking while demanding convergent judgment"
    p.font.italic = True
    p.font.size = Pt(13)
    
    # Slide 11: Type C Case Evidence
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Type C Case Evidence: P12"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "90% Authorship Despite 15-20% AI Contribution"
    
    p = tf.add_paragraph()
    p.text = "Five Iterative Cycles:"
    p.font.bold = True
    
    cycles = [
        '1. Initial vision: "Cursed princess in dragon form" (human)',
        '2. AI attempt: Opening lines about transformation',
        '3. Critical rejection: "I don\'t want repetitive use of \'whisper\'" (human judgment)',
        '4. AI refinement: Adjusted to "hidden hope"',
        '5. Narrative correction: "Princess seeks witch\'s help, not vice versa" (human agency)'
    ]
    
    for cycle in cycles:
        p = tf.add_paragraph()
        p.text = cycle
        p.level = 1
        p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = 'P16: "I asked, I instructed, I guided" → Language of choreography, not consumption'
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Slide 12: Change #6 - Limitations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Change #6 - Limitations Acknowledged"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What This Study Did Not Do"
    
    limitations = [
        "1. No L1 Comparison Group",
        "   - All 30 participants were L2 English learners",
        "   - Cannot isolate language background effects",
        "2. No Long-Period Observation",
        "   - Single 75-minute session per participant",
        "   - Cannot track development over time",
        "3. No No-AI Control Condition",
        "   - All participants used AI",
        "   - Cannot isolate pure AI effects",
        "4. Sample Size (N=30)",
        "   - Depth over breadth approach",
        "   - Rich qualitative data but limited generalizability"
    ]
    
    for lim in limitations:
        p = tf.add_paragraph()
        p.text = lim
        if lim[0].isdigit():
            p.level = 0
            p.font.bold = True
        else:
            p.level = 1
        p.font.size = Pt(13)
    
    # Slide 13: Future Study Plan
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Study Plan"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Addressing Current Limitations"
    
    p = tf.add_paragraph()
    p.text = "Design Changes:"
    p.font.bold = True
    
    changes = [
        "Groups: L2 only → L2 + L1 parallel groups",
        "Control: All AI-assisted → AI-assisted + No-AI condition",
        "Duration: Single 75-min session → Multi-session across semester",
        "Measures: Self-reported → + Computational attribution + Validated scales",
        "Analysis: Qualitative-dominant → Mixed: Thematic + Quantitative"
    ]
    
    for change in changes:
        p = tf.add_paragraph()
        p.text = change
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Key Additions: L1 control, No-AI baseline, Longitudinal design, Guilford's creativity scales, Pre-registration"
    p.font.italic = True
    p.font.size = Pt(13)
    
    # Slide 14: Theoretical Contribution Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Theoretical Contribution Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What This Revised Study Argues"
    
    p = tf.add_paragraph()
    p.text = "Core Claim:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Parameter manipulation is a pedagogical design lever that educators can use to dynamically adapt DDL scaffolding in AI-assisted creative writing"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Three Mechanisms Revealed:"
    p.font.bold = True
    
    mechanisms = [
        "1. Parameter → Interaction: Technical settings shape interaction types",
        "2. Interaction → Perception: Different types produce distinct authorship experiences",
        "3. Perception → Pedagogy: Learners intuitively detect and select scaffolding modes"
    ]
    
    for mech in mechanisms:
        p = tf.add_paragraph()
        p.text = mech
        p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Practical Implication: Educators need to understand parameter effects and adjust settings to match learner needs"
    p.font.italic = True
    p.font.size = Pt(13)
    
    # Slide 15: Response to Panel Feedback Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Response to Panel Feedback - Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "All Feedback Integrated"
    
    responses = [
        "Dr. Wang: DDL justification → AI as 'new DDL interface' argument",
        "Dr. Checketts: Poetry vs. machine logic → Heidegger-Strehovec dialectic",
        "Dr. Checketts & Dr. Hung: Creativity framework → Type C mechanism",
        "Dr. Harrington: Research questions → Unified parameter-interaction-perception RQ",
        "Dr. Hung: L2 framing → 'Novice writers' + limitations + future L1 groups"
    ]
    
    for resp in responses:
        p = tf.add_paragraph()
        p.text = resp
        p.level = 0
        p.font.size = Pt(14)
    
    # Slide 16: Open Questions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Open Questions for Discussion"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "For Panel Consideration"
    
    questions = [
        "1. Has the AI-as-DDL-interface argument successfully extended DDL theory?",
        "2. Is the three-interaction-types framework adequate without explicit Guilford measures?",
        "3. Does the 'novice writers' emphasis strengthen or weaken the study?",
        "4. To what extent do findings transfer to other creative genres or academic writing?",
        "5. What specific parameter configuration heuristics should educators follow?"
    ]
    
    for q in questions:
        p = tf.add_paragraph()
        p.text = q
        p.level = 0
        p.font.size = Pt(14)
    
    # Slide 17: Closing Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "DDL's Evolution, Not Death"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "The Original Challenge (Crosthwaite & Baisa 2023):"
    
    p = tf.add_paragraph()
    p.text = '"Generative AI represents the death of Data-driven Learning"'
    p.level = 1
    p.font.italic = True
    
    p = tf.add_paragraph()
    p.text = "Our Counter-Claim:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Generative AI represents DDL's evolution from static to dynamic"
    p.level = 1
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "From:"
    p.font.bold = True
    
    evolutions = [
        "Learner-initiated queries → System-initiated pattern presentation",
        "Fixed corpus evidence → Generative possibilities",
        "Manual concordance search → Conversational scaffolding",
        "One-size-fits-all → Parameter-adjustable adaptivity"
    ]
    
    for evo in evolutions:
        p = tf.add_paragraph()
        p.text = evo
        p.level = 1
        p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "The Future of DDL: Educators' capacity to exercise pedagogical judgment over AI configuration"
    p.font.italic = True
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Thank you. Questions welcome."
    p.font.bold = True
    p.font.size = Pt(16)
    
    # Save presentation
    output_path = "CRP_Second_Debate_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\n💡 Note: This script requires python-pptx library.")
        print("   Install it with: pip install python-pptx")
