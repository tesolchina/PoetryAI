"""
Generate PowerPoint slide for Oulipo and Heidegger Integration
Supplementary slides for CRP Second Debate
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_oulipo_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title - Theoretical Framework Addition
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Theoretical Framework Enhancement"
    subtitle.text = "Integrating Heidegger, Oulipo, and Strehovec\nAddressing Dr. Checketts' and Dr. Harrington's Feedback\n\nSupplementary Slides for CRP Second Debate"
    
    # Slide 2: The Three Contrasting Perspectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Three Contrasting Perspectives on Poetry & Machine Logic"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Addressing the Poetry-Machine Tension"
    
    p = tf.add_paragraph()
    p.text = "1. Heidegger (1984): Poetry as Antithesis to Machine Logic"
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    
    p = tf.add_paragraph()
    p.text = "• Poetry = purest human expression, incompatible with systematic procedures"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Categorical rejection: mechanizing what must remain ineffable"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "2. Oulipo (1960-present): Constraint as Creative Liberation"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = '• "Constraint liberates" - systematic procedures force unexplored linguistic territory'
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• 60-year track record: Perec's novel without letter 'e', N+7, mathematical structures"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "3. Strehovec (2023): Poetry as Cognitive Problem-Solving"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 192)
    
    p = tf.add_paragraph()
    p.text = "• Artistic creation = cognitive, goal-oriented activity"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "• Poetry composition as systematic craft involving learnable techniques"
    p.level = 1
    p.font.size = Pt(14)
    
    # Slide 3: Visual Comparison Table
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Theoretical Positions: Visual Comparison"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Create comparison text boxes
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(2.8)
    height = Inches(5.5)
    
    # Heidegger column
    box1 = slide.shapes.add_textbox(left, top, width, height)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Heidegger"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf1.add_paragraph()
    p.text = "\nPoetry vs. Machine Logic"
    p.font.size = Pt(14)
    p.font.bold = True
    
    content = [
        "↓",
        "Categorical separation",
        "Poetry = ineffable",
        "Science ≠ Art",
        "Machine = threat to authenticity"
    ]
    
    for item in content:
        p = tf1.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.alignment = PP_ALIGN.CENTER if item == "↓" else PP_ALIGN.LEFT
    
    # Oulipo column
    box2 = slide.shapes.add_textbox(left + Inches(3.3), top, width, height)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Oulipo"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf2.add_paragraph()
    p.text = "\nConstraint = Liberation"
    p.font.size = Pt(14)
    p.font.bold = True
    
    content = [
        "↓",
        "Deliberate constraints",
        "Systematic procedures",
        "Algorithmic generation",
        "Constraint → creativity"
    ]
    
    for item in content:
        p = tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.alignment = PP_ALIGN.CENTER if item == "↓" else PP_ALIGN.LEFT
    
    # Strehovec column
    box3 = slide.shapes.add_textbox(left + Inches(6.6), top, width, height)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "Strehovec"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 192)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf3.add_paragraph()
    p.text = "\nPoetry = Cognitive Work"
    p.font.size = Pt(14)
    p.font.bold = True
    
    content = [
        "↓",
        "Goal-oriented activity",
        "Learnable craft",
        "Problem-solving",
        "Methodology-supported"
    ]
    
    for item in content:
        p = tf3.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.alignment = PP_ALIGN.CENTER if item == "↓" else PP_ALIGN.LEFT
    
    # Slide 4: Oulipo Examples
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Oulipo: Historical Precedent for Constraint-Based Creativity"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "60 Years of Systematic Procedures in Literary Creation"
    
    p = tf.add_paragraph()
    p.text = "Key Examples:"
    p.font.bold = True
    
    examples = [
        "Georges Perec's La Disparition (1969): 300-page novel without letter 'e'",
        "Lipogram: Deliberately omitting certain letters throughout entire work",
        "N+7: Replace each noun with 7th noun following it in dictionary",
        "S+7: Same technique applied to other word classes",
        "Mathematical structures: Sonnets using Fibonacci sequences",
        "Constraint-based translation: Translating while maintaining original constraints"
    ]
    
    for example in examples:
        p = tf.add_paragraph()
        p.text = example
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = 'Core Principle: "La contrainte libère" (Constraint liberates)'
    p.font.italic = True
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # Slide 5: Our Synthesized Stance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Our Stance: Bounded Transparency in AI-Mediated Poetry"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Synthesizing Three Perspectives"
    
    p = tf.add_paragraph()
    p.text = "What We Accept:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    accepts = [
        "✓ From Oulipo: Systematic procedures can expand creative possibility",
        "✓ From Strehovec: Poetry composition as learnable, systematic craft",
        "✓ AI-generated patterns can function as productive constraints"
    ]
    
    for accept in accepts:
        p = tf.add_paragraph()
        p.text = accept
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "What We Reject:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)
    
    rejects = [
        "✗ Heidegger's categorical rejection of machine involvement",
        "✗ Treating AI as autonomous poet or creative agent",
        "✗ Passive acceptance of AI output without critical engagement"
    ]
    
    for reject in rejects:
        p = tf.add_paragraph()
        p.text = reject
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Core Position:"
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 192)
    
    p = tf.add_paragraph()
    p.text = "Machine influence is pedagogically valid when BOUNDED and TRANSPARENT"
    p.level = 1
    p.font.size = Pt(15)
    p.font.bold = True
    
    # Slide 6: AI as Pedagogical Constraint Provider
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI as Pedagogical Constraint Provider"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Not Creative Agent, But Oulipian Constraint Mechanism"
    
    p = tf.add_paragraph()
    p.text = "AI Functions Like Oulipian Constraints:"
    p.font.bold = True
    
    functions = [
        "Generate linguistic possibilities through systematic procedures",
        "Surface patterns that learners wouldn't discover alone",
        "Force exploration of unexplored linguistic territory",
        "Provide structured starting points for creative expansion"
    ]
    
    for func in functions:
        p = tf.add_paragraph()
        p.text = func
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Human Writers Maintain Critical Agency:"
    p.font.bold = True
    
    agency = [
        "Selection: Choose which AI suggestions to adopt",
        "Rejection: Discard suggestions that don't fit creative vision",
        "Revision: Modify AI output to match intent and voice",
        "Integration: Synthesize AI patterns into coherent creative work"
    ]
    
    for item in agency:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "AI = 'Potential Literature Machine' (Oulipo term)"
    p.font.italic = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # Slide 7: Oulipo vs AI Constraints Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Oulipo Constraints vs. AI Constraints"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Comparative Analysis"
    
    p = tf.add_paragraph()
    p.text = "Traditional Oulipo Constraints:"
    p.font.bold = True
    
    oulipo = [
        "Rigid: Lipogram completely excludes certain letters",
        "Static: N+7 follows fixed dictionary order",
        "Pre-defined: Mathematical structures set before writing",
        "Non-responsive: Constraint doesn't adapt to content"
    ]
    
    for item in oulipo:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "AI-Generated Constraints (Our Approach):"
    p.font.bold = True
    
    ai = [
        "Flexible: Suggestions can be partially adopted or modified",
        "Dynamic: AI responds to learner's developing text",
        "Context-aware: Suggestions adapt to theme, style, mood",
        "Responsive: Learners can request specific types of constraints"
    ]
    
    for item in ai:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "AI constraints are LESS rigid and MORE contextually responsive than Oulipian constraints"
    p.font.italic = True
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 192)
    
    # Slide 8: Theoretical Validation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why This Integration Strengthens Our Argument"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Three Forms of Validation"
    
    p = tf.add_paragraph()
    p.text = "1. Historical Precedent"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Oulipo's 60-year track record demonstrates systematic procedures CAN generate recognized literary achievements"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "2. Theoretical Grounding"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Bridges computational approach (Strehovec) with creative constraint tradition (Oulipo), addressing Heidegger's concerns"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "3. Pedagogical Justification"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "If strict Oulipian constraints produce creativity, then more flexible AI-generated constraints should enable authentic creative work"
    p.level = 1
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Key Insight: AI-assisted poetry aligns with 60-year tradition of constraint-based creation while introducing more adaptive, learner-responsive mechanisms"
    p.font.italic = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 128)
    
    # Slide 9: Addressing Panel Feedback
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How This Addresses Panel Feedback"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Direct Response to Dr. Checketts and Dr. Harrington"
    
    p = tf.add_paragraph()
    p.text = "Dr. Checketts' Concern:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Heidegger frames poetry as opposite of machine logic. Justify why poetry needs machine influence."'
    p.level = 1
    p.font.italic = True
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Our Response:"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Three-perspective framework explicitly engages Heidegger while showing Oulipo and Strehovec provide counterevidence that systematic procedures can serve poetry"
    p.level = 2
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "Dr. Harrington's Suggestion:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Explore Oulipo as relevant creative constraint tradition."'
    p.level = 1
    p.font.italic = True
    p.font.size = Pt(13)
    
    p = tf.add_paragraph()
    p.text = "Our Response:"
    p.level = 1
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Oulipo now serves as central theoretical bridge, providing both historical validation and conceptual vocabulary ('constraint liberates', 'potential literature machine')"
    p.level = 2
    p.font.size = Pt(13)
    
    # Slide 10: Integration with DDL Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Integration with DDL Framework"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "How Oulipo Enhances DDL Justification"
    
    p = tf.add_paragraph()
    p.text = "Traditional DDL:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Learner queries corpus → receives evidence → makes linguistic decisions"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Oulipian Constraint-Based Creation:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Writer imposes constraint → explores forced linguistic space → produces novel output"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "AI-Assisted DDL (Our Approach):"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "AI generates pattern constraints → learner notices/evaluates/revises → synthesizes with critical agency"
    p.level = 1
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "All three share: Active decision-making, Pattern-based learning, Constraint as productive force"
    p.font.italic = True
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    # Slide 11: Concluding Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Theoretical Contribution: Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What This Integration Achieves"
    
    achievements = [
        "Positions AI-assisted poetry within 60-year tradition of constraint-based creativity",
        "Addresses Heidegger's concerns while maintaining pedagogical validity",
        "Provides historical evidence that systematic procedures generate, not suppress, creativity",
        "Offers conceptual vocabulary from Oulipo ('constraint liberates', 'potential literature')",
        "Strengthens DDL justification by showing structural parallels with constraint-based learning",
        "Explicitly responds to both Dr. Checketts' and Dr. Harrington's feedback"
    ]
    
    for achievement in achievements:
        p = tf.add_paragraph()
        p.text = achievement
        p.level = 0
        p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = ""
    
    p = tf.add_paragraph()
    p.text = "AI as bounded, transparent constraint provider = pedagogically defensible evolution of both DDL and Oulipian traditions"
    p.font.italic = True
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(0, 0, 192)
    
    # Save presentation
    output_path = "CRP_Oulipo_Integration_Slides.pptx"
    prs.save(output_path)
    print(f"✅ Oulipo integration slides created successfully: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    try:
        create_oulipo_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\n💡 Note: This script requires python-pptx library.")
        print("   Install it with: pip install python-pptx")
