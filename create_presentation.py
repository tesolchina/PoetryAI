from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 85)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 45, 85)
    
    # Add horizontal line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(25, 45, 85)
    line.line.width = Pt(2)
    
    # Add content
    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(5.3))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# Slide 1: Title Slide
add_title_slide(prs, 
    "Parameter Configuration & L2 Creative Writing",
    "How AI Temperature Settings Shape DDL Scaffolding in Poetry")

# Slide 2: Research Problem & Design
add_content_slide(prs,
    "Research Question & Experimental Design",
    [
        "🔍 Does parameter manipulation enable dynamic DDL scaffolding?",
        "",
        "📊 2×2 Experimental Design (N=10):",
        "  • Structured parameters (Temp 0.3, Top-p 0.4) vs.",
        "  • Exploratory parameters (Temp 0.8, Top-p 0.9)",
        "  • Awareness condition (aware vs. unaware)",
        "",
        "Three Interaction Types:",
        "  • Type A: Constraint Repair (corrective feedback)",
        "  • Type B: Exemplar Giving (preset options)",
        "  • Type C: Surprise Harvest (unexpected ideas)"
    ])

# Slide 3: Key Findings
add_content_slide(prs,
    "Major Findings",
    [
        "📈 7× Increase in Type C Interactions:",
        "  • Structured: 5% Type C, 60% Type A, 35% Type B",
        "  • Exploratory: 35% Type C, 20% Type A, 45% Type B",
        "  • χ² = 24.3, p < .001 (highly significant)",
        "",
        "⚠️  The \"Helpful but Alienating\" Paradox:",
        "  • 78% rated Type B as \"most helpful\"",
        "  • But Type B negatively correlated with authorship",
        "  • Structured rooms: avg 28% self-authorship",
        "  • Exploratory rooms: avg 48% self-authorship",
        "",
        "✨ Type C Enables Ownership:",
        "  • Participant preferring Type C reported 80% authorship"
    ])

# Slide 4: Implications & Conclusion
add_content_slide(prs,
    "Implications & Conclusions",
    [
        "🎯 Parameter Literacy as Essential Competence",
        "  Educators must understand how temperature shapes pedagogy",
        "",
        "🔄 Parameter Configuration as Pedagogical Lever",
        "  AI evolves from pattern-enforcer to pattern-extender",
        "",
        "📚 DDL's Evolution, Not Death",
        "  AI preserves constructivist discovery when configured for",
        "  Type C interactions (exploratory parameters)",
        "",
        "💡 Key Insight:",
        "  The difference between alienation and authorship may hinge",
        "  on adjusting a temperature setting from 0.3 to 0.8"
    ])

# Slide 5: Platform Introduction
add_content_slide(prs,
    "The Research Platform: poetry.aitutor.ink",
    [
        "🏗️ Purpose-Built Infrastructure",
        "  • Custom poetry writing platform (not commercial AI tool)",
        "  • Transparent experimental laboratory with full parameter control",
        "",
        "🎨 Design Principles",
        "  • Guide, not ghostwriter: AI as writing coach & collaborator",
        "  • Parameters exposed as pedagogical variables (not hidden)",
        "  • Prompt-designed to support all 3 interaction types",
        "  • Response limit: 40-80 words (prevent overwhelming)",
        "",
        "🔬 Research Features",
        "  • Complete interaction logging: every keystroke, revision, suggestion",
        "  • 4 Virtual rooms with different parameter configurations",
        "  • Identical prompts, only parameters vary (clean experiment)",
        "  • Built with Python/Flask + OpenRouter API + Claude Sonnet 4"
    ])

# Save presentation
output_path = r"c:\Users\ruobin Yu\.vscode\PoetryAI-6\DDL_Special_Issue_Presentation_5slides.pptx"
prs.save(output_path)
print(f"✓ Presentation created: {output_path}")
