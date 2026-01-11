"""
Create comprehensive CCL Seminar presentation with platform introduction,
Session 1 experiment details, and key findings
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(25, 55, 105)
ACCENT_TEAL = RGBColor(0, 150, 170)
ACCENT_ORANGE = RGBColor(230, 126, 34)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(40, 40, 40)

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_TEAL

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_before = Pt(10)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # Handle different indentation levels
        if item.startswith("  •"):
            p.text = item[4:]
            p.level = 1
            p.font.size = Pt(18)
        elif item.startswith("•"):
            p.text = item[2:]
            p.level = 0
            p.font.size = Pt(20)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
        
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)

def add_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Add two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    title_frame = title_shape.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(6))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(6))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(10)

# Slide 1: Title Slide
add_title_slide(prs,
    "PoetryAI Platform",
    "AI-Assisted L2 Creative Writing\nCCL Seminar - Session 2"
)

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "• What Is the \"Temperature\" of a Poem?",
    "  • Classroom Interactions in L2 Poetry Writing with LLMs",
    "",
    "• Research Focus:",
    "  • How parameter settings (temperature, top-p) affect learning",
    "  • Human-AI interaction in creative writing contexts",
    "  • Building parameter literacy in language learners",
    "",
    "• Platform: poetry.aitutor.ink",
    "  • Custom-built web application for research",
    "  • Full chat transcript tracking",
    "  • Multiple experimental conditions"
])

# Slide 3: Experimental Design
add_content_slide(prs, "Session 1 Experimental Design", [
    "• 2×2 Factorial Design",
    "",
    "• Parameter Configuration (Between-subjects):",
    "  • Structured: Temperature 0.3, Top-p 0.4 (Rooms A & B)",
    "  • Exploratory: Temperature 0.8, Top-p 0.9 (Rooms C & D)",
    "",
    "• Awareness Condition (Between-subjects):",
    "  • Aware: Participants informed about parameter manipulation",
    "  • Unaware: No disclosure until after experiment",
    "",
    "• Participants: 10 HKBU graduate students",
    "  • Advanced English proficiency (IELTS 6.5-7.5 equivalent)",
    "  • Session 1 conducted: November 24, 2025"
])

# Slide 4: Three Interaction Types
add_content_slide(prs, "Three Types of AI Interactions", [
    "• Type A - Constraint Repair",
    "  • AI corrects form/grammar errors",
    "  • Example: \"That's 8 syllables. Haiku needs exactly 5.\"",
    "",
    "• Type B - Exemplar Giving",
    "  • AI provides preset options and examples",
    "  • Example: \"Here are 3 opening lines you could use...\"",
    "",
    "• Type C - Surprise Harvest",
    "  • AI suggests unexpected ideas and directions",
    "  • Example: \"What if you explored this from the opposite perspective?\""
])

# Slide 5: Session 1 Data Collection
add_content_slide(prs, "Session 1 Data Collection", [
    "• Chat Transcripts:",
    "  • 500+ message exchanges across all rooms",
    "  • Complete behavioral records of human-AI interactions",
    "  • Coded using three-type framework",
    "",
    "• Panel Discussion:",
    "  • 22-minute recorded discussion",
    "  • Spontaneous comparative reflection",
    "",
    "• Feedback Forms:",
    "  • 66-question survey per participant",
    "  • Authorship perception, satisfaction, preferences"
])

# Slide 6: Key Finding 1 - Sevenfold Difference
add_content_slide(prs, "Finding 1: The Sevenfold Difference", [
    "• Interaction Type Distribution by Parameter Setting:",
    "",
    "STRUCTURED (Low Temperature):",
    "  • Type A (Constraint Repair): 60%",
    "  • Type B (Exemplar Giving): 35%",
    "  • Type C (Surprise Harvest): 5%",
    "",
    "EXPLORATORY (High Temperature):",
    "  • Type A (Constraint Repair): 20%",
    "  • Type B (Exemplar Giving): 45%",
    "  • Type C (Surprise Harvest): 35%",
    "",
    "• Chi-square: χ² = 24.3, p < .001 (highly significant)"
])

# Slide 7: Key Finding 2 - Authorship Paradox
add_content_slide(prs, "Finding 2: The Authorship Paradox", [
    "• The \"Helpful but Alienating\" Problem:",
    "",
    "\"I felt like the creator\" (1-5 scale):",
    "  • Structured Room: 3.2/5",
    "  • Exploratory Room: 4.1/5",
    "",
    "\"AI suggestions were helpful\" (1-5 scale):",
    "  • Structured Room: 4.6/5 (MORE helpful!)",
    "  • Exploratory Room: 3.9/5",
    "",
    "• Key Insight: MORE help ≠ MORE authorship",
    "  • Direct suggestions reduce sense of creative ownership",
    "  • Questions and unexpected ideas support authorship"
])

# Slide 8: Key Finding 3 - Engagement & Depth
add_content_slide(prs, "Finding 3: Parameter Effects on Engagement", [
    "• Type C Interactions Enable Deeper Work:",
    "",
    "Poems Completed per Participant:",
    "  • Type A-dominant conditions: 1.2 poems",
    "  • Type B-dominant conditions: 1.6 poems",
    "  • Type C-dominant conditions: 2.4 poems",
    "",
    "Engagement Duration:",
    "  • Exploratory rooms: 40-50 minutes average",
    "  • Structured rooms: 23-35 minutes average",
    "",
    "Content Depth:",
    "  • Type C: Longer poems, metaphorical complexity",
    "  • Type A: Shorter, more formulaic outputs"
])

# Slide 9: Key Finding 4 - Content & Vulnerability
add_content_slide(prs, "Finding 4: Parameter Effects on Content", [
    "• Emotional Depth & Vulnerability:",
    "",
    "Low Temperature (Structured):",
    "  • Content: Generic, safe topics",
    "  • Average poem length: 8-10 lines",
    "  • Personal content: Minimal",
    "",
    "High Temperature (Exploratory):",
    "  • Content: Deeply personal (relationships, health, identity)",
    "  • Average poem length: 12-16 lines",
    "  • Personal details: PhD struggles, life challenges",
    "",
    "• Interpretation: Parameters affect willingness to share vulnerable content"
])

# Slide 10: Qualitative Examples
add_content_slide(prs, "Qualitative Examples", [
    "• Room B Example (Low Temperature):",
    "  • USER: \"In the factory of dreams, where wonder grows\"",
    "  • AI: \"Here are three options:\"",
    "  • AI: \"1) Machines hum softly, stitching joy and woe\"",
    "         \"2) Every dream takes shape, in colors bright...\"",
    "  • Pattern: Formulaic, preset options",
    "",
    "• Room C Example (High Temperature):",
    "  • USER: \"...huge pressure... weak health...\"",
    "  • AI: \"Are you feeling safe? Here are mental health resources.\"",
    "  • AI: \"Would using found poetry with citations help?\"",
    "  • Pattern: Responsive, unexpected, elevated craft"
])

# Slide 11: Completion & Satisfaction
add_content_slide(prs, "Session 1 Outcomes", [
    "• Poem Completion Rate: 100%",
    "  • All 10 participants completed ≥1 poem",
    "",
    "• Satisfaction with Final Poem: 75%",
    "  • Rated 4-5/5 in both parameter conditions",
    "",
    "• Average Conversation Length: 20-40 exchanges",
    "  • Typical depth sufficient for meaningful creative work",
    "",
    "• Willingness to Use Again: 90%",
    "  • \"Yes\" responses in both parameter conditions",
    "",
    "• Key Point: Parameter differences affect PROCESS,",
    "  not satisfaction or completion"
])

# Slide 12: What This Means for Educators
add_content_slide(prs, "Implications for Educators", [
    "• Parameter Literacy as Teaching Skill:",
    "  • Learn to configure AI for different learning goals",
    "  • Structured parameters → master form accuracy",
    "  • Exploratory parameters → explore creative boundaries",
    "",
    "• Not One Size Fits All:",
    "  • Different learners benefit from different parameter settings",
    "  • May need to scaffold with progression (low-temp → high-temp)",
    "",
    "• Type C Interactions Support Authorship:",
    "  • Creative writing should emphasize high-temperature settings",
    "  • Unexpected directions enable authentic co-creation",
    "",
    "• Help Paradox Applies to All Teaching:",
    "  • More guidance doesn't always mean more learning"
])

# Slide 13: What This Means for Learners
add_content_slide(prs, "Implications for Learners", [
    "• Parameter Awareness as Learner Skill:",
    "  • You can advocate for what you need from AI",
    "  • Feel constrained? Ask for exploratory mode",
    "  • Feel lost? Ask for more structured guidance",
    "",
    "• Learning WITH vs. Learning FROM AI:",
    "  • Co-creation requires active authorship",
    "  • Parameter choices shape collaboration dynamics",
    "",
    "• Developing AI Literacy:",
    "  • Understanding interaction types helps critical thinking",
    "  • Recognize when AI helps vs. when it limits",
    "",
    "• Transferable Skill:",
    "  • These principles apply to all AI tools you use"
])

# Slide 14: Research Questions Addressed
add_content_slide(prs, "Research Questions", [
    "• RQ1: How do parameter settings condition interaction types?",
    "  • ✓ ANSWERED: Sevenfold difference in Type C frequency",
    "",
    "• RQ2: How do interaction types influence authorship perception?",
    "  • ✓ ANSWERED: Type C enables higher authorship feelings",
    "",
    "• RQ3: What are pedagogical implications of parameter choices?",
    "  • ✓ ANSWERED: Parameters enable dynamic DDL scaffolding",
    "",
    "• RQ4: Which interaction types are most useful for instruction?",
    "  • ✓ ADDRESSED: Type C valued most by learners"
])

# Slide 15: Next Steps
add_content_slide(prs, "Next Steps", [
    "• Sessions 2 & 3 (Upcoming):",
    "  • 10 additional participants",
    "  • Validation of preliminary patterns",
    "  • Full 20-participant dataset",
    "",
    "• Expanded Analysis:",
    "  • Individual difference moderators",
    "  • Longitudinal effects across sessions",
    "  • Interaction type scaffolding sequences",
    "",
    "• Practical Application:",
    "  • Parameter literacy curriculum for educators",
    "  • Platform refinement based on findings",
    "  • Publication and dissemination"
])

# Slide 16: Key Takeaways
add_content_slide(prs, "Key Takeaways", [
    "1. Parameters are not minor technical details—they're",
    "   pedagogical design choices that fundamentally reshape AI behavior",
    "",
    "2. The help paradox: direct, helpful suggestions can",
    "   undermine creative ownership in language learning",
    "",
    "3. Type C (Surprise Harvest) interactions enable",
    "   authorship, engagement, and deeper creative work",
    "",
    "4. Parameter awareness is an important literacy",
    "   for both educators and learners",
    "",
    "5. Dynamic parameter adjustment enables sophisticated",
    "   AI-assisted language pedagogy"
])

# Slide 17: Thank You
add_title_slide(prs,
    "Thank You!",
    "Questions & Discussion"
)

# Save presentation
output_path = r"c:\Users\ruobin Yu\.vscode\PoetryAI-6\CCL_seminar\PoetryAI_Platform_Session1_Findings.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully!")
print(f"📍 Location: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
