from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = LIGHT_BLUE

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = DARK_BLUE
    title_shape.line.color.rgb = DARK_BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_frame.margin_left = Inches(0.5)
    title_frame.vertical_anchor = 1
    
    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, content in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = content
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_quote_slide(prs, quote, author):
    """Add a quote slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BLUE
    
    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
    quote_frame = quote_box.text_frame
    quote_frame.word_wrap = True
    quote_p = quote_frame.paragraphs[0]
    quote_p.text = f'"{quote}"'
    quote_p.font.size = Pt(28)
    quote_p.font.italic = True
    quote_p.font.color.rgb = WHITE
    quote_p.alignment = PP_ALIGN.CENTER
    
    # Author
    author_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.8))
    author_frame = author_box.text_frame
    author_p = author_frame.paragraphs[0]
    author_p.text = f"— {author}"
    author_p.font.size = Pt(20)
    author_p.font.color.rgb = WHITE
    author_p.alignment = PP_ALIGN.CENTER

# Slide 1: Title Slide
add_title_slide(prs, "AI and Creative Poetry Writing", 
                "Exploring Human-AI Collaboration in Language Learning\nCCL Seminar")

# Slide 2: The Story
add_content_slide(prs, "The Story of Charles Bernstein", [
    "• 2019: Publication of Poetry Has No Future Unless It Comes to an End",
    "• Co-created with AI trained on Bernstein's own poetry",
    "• The AI called itself the 'Synthetic Brother'",
    "• Yet... the AI remained uncredited in the author column",
    "• This erasure reveals our ambivalence about AI creativity"
])

# Slide 3: The Central Questions
add_quote_slide(prs, "In an age of technological advancement, where exactly do we humans stand?", 
                "Central Question")

# Slide 4: The Challenge
add_content_slide(prs, "Our Era of Ambivalence", [
    "• AI brings exciting new possibilities to creative fields",
    "• Yet we face urgent questions:",
    "   - What happens to human originality?",
    "   - What does authorship mean anymore?",
    "   - Can machines replace human creativity?"
])

# Slide 5: Skepticism
add_content_slide(prs, "The Academic Skepticism", [
    "• AI-generated poetry often lacks genuine emotion",
    "• Missing contextual depth and nuanced understanding",
    "• Human bias against machine-created work",
    "• Fear of losing human creativity: 'anxiety of machine influence'",
    "• Yet... this skepticism is beginning to shift"
])

# Slide 6: The Shift
add_quote_slide(prs, "What if AI isn't a threat, but a catalyst for creativity?", 
                "The Emerging Perspective")

# Slide 7: Why Poetry?
add_content_slide(prs, "Why Poetry (and Not Prose)?", [
    "• Poetry tolerates ambiguity and formal flexibility",
    "• Grammar and syntax are not strictly enforced",
    "• Readers expect uncertainty and abstraction",
    "• Lower technical barrier than fiction or prose",
    "• BUT: Technical possibility ≠ Capturing true poetry"
])

# Slide 8: What AI Offers
add_content_slide(prs, "What Can AI Bring to Poetry?", [
    "1. Idea Prompt & Diversity Explorer",
    "   - Draw from thousands of poems across centuries",
    "   - Fresh perspectives unconstrained by human bias",
    "",
    "2. Democratization & Personalization",
    "   - Make poetry creation accessible to everyone",
    "",
    "3. Mastery Imitator",
    "   - Simulate styles and approaches of great poets"
])

# Slide 9: The Critical Issue
add_content_slide(prs, "The Power of Human Decisions", [
    "• A poem's power emerges from interplay, not single elements",
    "• Nuance lives in how poets make choices",
    "• Selecting one word from many synonyms",
    "• Embracing unexpected AI suggestions creatively",
    "• These micro-decisions are where poetry becomes poetry",
    "• These decisions must be documented and studied"
])

# Slide 10: The Irreducible
add_quote_slide(prs, "The incommunicable, deeply private aspects of human expression cannot be reduced to algorithms.", 
                "The Irreducible Humanity")

# Slide 11: The Path Forward
add_content_slide(prs, "The Path Forward: Collaboration", [
    "• Human-Machine In-A-Loop methodology",
    "• Genuine partnership, not human editing machine output",
    "• Humans continuously add interpretive layers",
    "• Requires Self-Efficacy:",
    "   - Participants co-create, not just edit",
    "• Requires Human-Centered Mindset:",
    "   - AI must serve human development"
])

# Slide 12: Questions to Guide Us
add_content_slide(prs, "Questions Before Us", [
    "• In an era when machines generate poetry in seconds,",
    "  what is the irreducible value of human creativity?",
    "",
    "• What will poetry become when humans and machines",
    "  collaborate thoughtfully and intentionally?",
    "",
    "• How can we design technologies that strengthen",
    "  human capacity and amplify human voice?"
])

# Slide 13: Closing
add_title_slide(prs, "The Answer Lies in Partnership", 
                "Not choosing between human or machine,\nbut understanding how to strengthen both")

# Save presentation
output_path = r"c:\Users\ruobin Yu\.vscode\PoetryAI-6\CCL_seminar\CCL_Seminar_AI_Poetry_Introduction.pptx"
prs.save(output_path)
print(f"Presentation created successfully at: {output_path}")
